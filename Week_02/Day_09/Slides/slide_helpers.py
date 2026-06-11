"""python-pptx slide builders for Day 9 deck.

All slide layouts are constructed from blank layouts to keep full control over
positioning and theming. Every content slide gets:
  - title bar with accent underline
  - body region with bullets / table / image / chart
  - footer with section name + source URL
  - speaker notes (the trainer "script")
"""
from __future__ import annotations
from dataclasses import dataclass, field
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

import theme as T


# --------------------------------------------------------------------------- #
# Presentation factory
# --------------------------------------------------------------------------- #
def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = T.SLIDE_W
    prs.slide_height = T.SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _paint_bg(slide, T.BG_DARK)
    return slide


# --------------------------------------------------------------------------- #
# Low-level primitives
# --------------------------------------------------------------------------- #
def _paint_bg(slide, color: RGBColor):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, T.SLIDE_W, T.SLIDE_H
    )
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    # send to back by moving the shape's XML element to first position
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)
    return rect


def add_rect(slide, x, y, w, h, fill: RGBColor | None = None,
             line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text: str, *,
             font=T.FONT_BODY, size=18, bold=False, italic=False,
             color: RGBColor = T.FG_TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets: Sequence[str], *,
                size=20, color: RGBColor = T.FG_TEXT,
                bullet_color: RGBColor = T.ACCENT_TEAL,
                line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_top = Pt(2)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        # bullet glyph
        r1 = p.add_run()
        r1.text = "▸ "
        r1.font.name = T.FONT_BODY
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        # body text (supports "Lead: rest" emphasis)
        if ":" in item and item.index(":") < 40:
            lead, rest = item.split(":", 1)
            r2 = p.add_run()
            r2.text = lead + ":"
            r2.font.name = T.FONT_BODY
            r2.font.size = Pt(size)
            r2.font.bold = True
            r2.font.color.rgb = T.ACCENT_YELLOW
            r3 = p.add_run()
            r3.text = rest
            r3.font.name = T.FONT_BODY
            r3.font.size = Pt(size)
            r3.font.color.rgb = color
        else:
            r2 = p.add_run()
            r2.text = item
            r2.font.name = T.FONT_BODY
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
    return tb


# --------------------------------------------------------------------------- #
# Slide chrome (title bar + footer)
# --------------------------------------------------------------------------- #
def add_title_bar(slide, title: str, *, section: str | None = None,
                  accent: RGBColor = T.ACCENT_TEAL):
    # accent stripe on left
    add_rect(slide, 0, 0, Inches(0.12), T.SLIDE_H, fill=accent)
    # section eyebrow
    if section:
        add_text(slide, T.MARGIN_X, Inches(0.25), Inches(8), Inches(0.3),
                 section.upper(), size=11, bold=True, color=accent,
                 font=T.FONT_BODY)
    # title text
    add_text(slide, T.MARGIN_X, Inches(0.55), Inches(12.3), Inches(0.9),
             title, size=30, bold=True, color=T.FG_TEXT, font=T.FONT_TITLE)
    # underline
    add_rect(slide, T.MARGIN_X, Inches(1.45), Inches(1.4), Inches(0.05),
             fill=accent)


def add_footer(slide, *, page: int, total: int, source: str | None = None):
    # Source on left
    if source:
        add_text(slide, T.MARGIN_X, Inches(7.05), Inches(10),
                 Inches(0.35), f"Source: {source}",
                 size=9, color=T.FG_DIM, italic=True)
    # Page number on right
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.6),
             Inches(0.35), f"{page} / {total}",
             size=10, color=T.FG_MUTED, align=PP_ALIGN.RIGHT)


def add_notes(slide, script: str):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = script


# --------------------------------------------------------------------------- #
# Compound slide types
# --------------------------------------------------------------------------- #
def title_slide(prs, *, title: str, subtitle: str, presenter: str = "Medallion Pipeline Bootcamp"):
    slide = blank_slide(prs)
    # big accent block
    add_rect(slide, 0, 0, Inches(4.5), T.SLIDE_H, fill=T.BG_PANEL)
    add_rect(slide, Inches(4.5), 0, Inches(0.08), T.SLIDE_H, fill=T.ACCENT_TEAL)
    # eyebrow
    add_text(slide, Inches(5.0), Inches(1.8), Inches(8), Inches(0.4),
             "DAY 9", size=14, bold=True, color=T.ACCENT_TEAL)
    # main title
    add_text(slide, Inches(5.0), Inches(2.2), Inches(8), Inches(1.6),
             title, size=44, bold=True, color=T.FG_TEXT, font=T.FONT_TITLE)
    # subtitle
    add_text(slide, Inches(5.0), Inches(4.0), Inches(8), Inches(1.2),
             subtitle, size=20, color=T.FG_MUTED, italic=True)
    # divider
    add_rect(slide, Inches(5.0), Inches(5.5), Inches(2), Inches(0.04),
             fill=T.ACCENT_YELLOW)
    add_text(slide, Inches(5.0), Inches(5.7), Inches(8), Inches(0.4),
             presenter, size=14, color=T.FG_TEXT)
    add_text(slide, Inches(5.0), Inches(6.1), Inches(8), Inches(0.4),
             "Monitoring · Logging · Data Quality · Scaling",
             size=12, color=T.FG_DIM)
    # corner decoration
    add_text(slide, Inches(0.8), Inches(6.6), Inches(3.5), Inches(0.6),
             "Theory Block · 3–4 hours", size=13, color=T.ACCENT_YELLOW, bold=True)
    return slide


def section_divider(prs, *, number: int, title: str, summary: str,
                    accent: RGBColor = T.ACCENT_TEAL):
    slide = blank_slide(prs)
    # giant number
    add_text(slide, Inches(0.6), Inches(1.4), Inches(4), Inches(4.5),
             f"{number:02d}", size=220, bold=True, color=accent,
             font=T.FONT_TITLE)
    add_text(slide, Inches(4.8), Inches(2.4), Inches(8), Inches(0.4),
             "SECTION", size=14, bold=True, color=accent)
    add_text(slide, Inches(4.8), Inches(2.8), Inches(8), Inches(1.2),
             title, size=40, bold=True, color=T.FG_TEXT, font=T.FONT_TITLE)
    add_rect(slide, Inches(4.8), Inches(4.1), Inches(1.5), Inches(0.05),
             fill=accent)
    add_text(slide, Inches(4.8), Inches(4.3), Inches(8), Inches(2.5),
             summary, size=18, color=T.FG_MUTED, italic=True)
    return slide


def content_slide(prs, *, title: str, section: str, bullets: Sequence[str],
                  source: str | None = None, accent=T.ACCENT_TEAL,
                  page: int = 0, total: int = 0, notes: str = "",
                  right_panel=None):
    """Content slide with bullets on the left and an optional right panel.

    right_panel can be:
      - ('text', str)         : multiline supporting paragraph
      - ('image', path)       : an image file path
      - ('shapes', callable)  : callable(slide, x, y, w, h) that draws shapes
    """
    slide = blank_slide(prs)
    add_title_bar(slide, title, section=section, accent=accent)

    body_y = Inches(1.7)
    body_h = Inches(5.2)
    if right_panel is None:
        add_bullets(slide, T.MARGIN_X, body_y, Inches(12.3), body_h,
                    bullets, size=20)
    else:
        # left bullets (≈55% width), right panel (≈40% width)
        add_bullets(slide, T.MARGIN_X, body_y, Inches(7.2), body_h,
                    bullets, size=18)
        rx, ry, rw, rh = Inches(7.9), body_y, Inches(5.0), body_h
        kind, payload = right_panel
        if kind == "text":
            add_rect(slide, rx, ry, rw, rh, fill=T.BG_PANEL)
            add_text(slide, rx + Inches(0.2), ry + Inches(0.2),
                     rw - Inches(0.4), rh - Inches(0.4),
                     payload, size=14, color=T.FG_MUTED, italic=True)
        elif kind == "image":
            try:
                slide.shapes.add_picture(payload, rx, ry,
                                         width=rw, height=rh)
            except Exception:
                add_rect(slide, rx, ry, rw, rh, fill=T.BG_PANEL)
                add_text(slide, rx + Inches(0.2), ry + Inches(0.2),
                         rw - Inches(0.4), rh - Inches(0.4),
                         f"[image missing: {payload}]",
                         size=12, color=T.FG_DIM)
        elif kind == "shapes":
            payload(slide, rx, ry, rw, rh)

    add_footer(slide, page=page, total=total, source=source)
    add_notes(slide, notes)
    return slide


def table_slide(prs, *, title: str, section: str, headers: Sequence[str],
                rows: Sequence[Sequence[str]], source: str | None = None,
                accent=T.ACCENT_TEAL, page=0, total=0, notes="",
                first_col_accent: bool = True):
    slide = blank_slide(prs)
    add_title_bar(slide, title, section=section, accent=accent)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_x, tbl_y = T.MARGIN_X, Inches(1.75)
    tbl_w, tbl_h = Inches(12.3), Inches(5.0)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, tbl_x, tbl_y,
                                       tbl_w, tbl_h)
    tbl = tbl_shape.table

    # column widths — first column wider when there are <=4 cols
    if n_cols <= 4:
        first_w = int(tbl_w * 0.35)
        rest = (tbl_w - first_w) // (n_cols - 1)
        tbl.columns[0].width = first_w
        for c in range(1, n_cols):
            tbl.columns[c].width = rest
    else:
        for c in range(n_cols):
            tbl.columns[c].width = tbl_w // n_cols

    # header row
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        _set_cell_text(cell, h, size=14, bold=True,
                       color=T.BG_DARK, align=PP_ALIGN.LEFT)

    # body rows
    for r_idx, row in enumerate(rows, start=1):
        row_color = T.BG_PANEL if r_idx % 2 == 1 else T.BG_PANEL_2
        for c, value in enumerate(row):
            cell = tbl.cell(r_idx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            is_lead = (c == 0 and first_col_accent)
            _set_cell_text(cell, value,
                           size=12,
                           bold=is_lead,
                           color=T.ACCENT_YELLOW if is_lead else T.FG_TEXT,
                           align=PP_ALIGN.LEFT)

    add_footer(slide, page=page, total=total, source=source)
    add_notes(slide, notes)
    return slide


def _set_cell_text(cell, text, *, size=12, bold=False,
                   color: RGBColor = T.FG_TEXT, align=PP_ALIGN.LEFT):
    cell.margin_left = Pt(6)
    cell.margin_right = Pt(6)
    cell.margin_top = Pt(4)
    cell.margin_bottom = Pt(4)
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # clear default text
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = T.FONT_BODY
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def chart_slide(prs, *, title: str, section: str, chart_path: str,
                caption: str = "", source: str | None = None,
                accent=T.ACCENT_TEAL, page=0, total=0, notes=""):
    slide = blank_slide(prs)
    add_title_bar(slide, title, section=section, accent=accent)
    img_x, img_y = T.MARGIN_X, Inches(1.7)
    img_w, img_h = Inches(12.3), Inches(4.8)
    try:
        slide.shapes.add_picture(chart_path, img_x, img_y,
                                 width=img_w, height=img_h)
    except Exception:
        add_rect(slide, img_x, img_y, img_w, img_h, fill=T.BG_PANEL)
        add_text(slide, img_x + Inches(0.3), img_y + Inches(0.3),
                 img_w, img_h, f"[chart missing: {chart_path}]",
                 size=14, color=T.FG_DIM)
    if caption:
        add_text(slide, T.MARGIN_X, Inches(6.6), Inches(12.3),
                 Inches(0.4), caption, size=12, color=T.FG_MUTED,
                 italic=True, align=PP_ALIGN.CENTER)
    add_footer(slide, page=page, total=total, source=source)
    add_notes(slide, notes)
    return slide


def code_slide(prs, *, title: str, section: str, code: str, language: str = "",
               caption: str = "", source: str | None = None,
               accent=T.ACCENT_TEAL, page=0, total=0, notes=""):
    slide = blank_slide(prs)
    add_title_bar(slide, title, section=section, accent=accent)
    # code panel
    cx, cy = T.MARGIN_X, Inches(1.75)
    cw, ch = Inches(12.3), Inches(4.9)
    add_rect(slide, cx, cy, cw, ch, fill=T.BG_PANEL)
    if language:
        add_text(slide, cx + Inches(0.2), cy + Inches(0.05),
                 Inches(4), Inches(0.3), language,
                 size=10, color=T.ACCENT_YELLOW, italic=True, bold=True)
    tb = slide.shapes.add_textbox(cx + Inches(0.25),
                                  cy + Inches(0.35),
                                  cw - Inches(0.5),
                                  ch - Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.1
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = T.FONT_MONO
        r.font.size = Pt(13)
        r.font.color.rgb = T.FG_TEXT
    if caption:
        add_text(slide, T.MARGIN_X, Inches(6.7), Inches(12.3),
                 Inches(0.3), caption, size=12, color=T.FG_MUTED,
                 italic=True)
    add_footer(slide, page=page, total=total, source=source)
    add_notes(slide, notes)
    return slide


def quiz_slide(prs, *, title: str, section: str, question: str,
               qtype: str, options: Sequence[str] | None = None,
               accent=T.ACCENT_YELLOW, page=0, total=0, notes=""):
    """qtype in {MCQ, MULTI, MATCH, FILL, BEST, TF}."""
    slide = blank_slide(prs)
    add_title_bar(slide, title, section=section, accent=accent)
    # quiz badge
    badge_w = Inches(2.2)
    add_rect(slide, T.MARGIN_X, Inches(1.65), badge_w, Inches(0.42),
             fill=accent)
    badge_label = {
        "MCQ": "QUIZ · MCQ",
        "MULTI": "QUIZ · MULTI-SELECT",
        "MATCH": "QUIZ · MATCH THE FOLLOWING",
        "FILL": "QUIZ · FILL IN THE BLANKS",
        "BEST": "QUIZ · CHOOSE THE BEST",
        "TF": "QUIZ · TRUE / FALSE",
    }.get(qtype, "QUIZ")
    add_text(slide, T.MARGIN_X, Inches(1.66), badge_w, Inches(0.4),
             badge_label, size=12, bold=True, color=T.BG_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, T.MARGIN_X, Inches(2.25), Inches(12.3),
             Inches(1.3), question, size=22, bold=True,
             color=T.FG_TEXT)

    if options:
        prefixes = ["A.", "B.", "C.", "D.", "E.", "F."]
        col_x = T.MARGIN_X
        col_w = Inches(12.3)
        opt_y = Inches(3.7)
        for i, opt in enumerate(options):
            y = opt_y + Inches(0.7 * i)
            add_rect(slide, col_x, y, Inches(0.55), Inches(0.55),
                     fill=T.BG_PANEL_2)
            add_text(slide, col_x, y, Inches(0.55), Inches(0.55),
                     prefixes[i], size=16, bold=True,
                     color=accent, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_rect(slide, col_x + Inches(0.65), y, col_w - Inches(0.65),
                     Inches(0.55), fill=T.BG_PANEL)
            add_text(slide, col_x + Inches(0.85), y, col_w - Inches(0.85),
                     Inches(0.55), opt, size=15, color=T.FG_TEXT,
                     anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page=page, total=total)
    add_notes(slide, notes)
    return slide


def answer_slide(prs, *, title: str, section: str, answer: str,
                 explanation: str, accent=T.ACCENT_GREEN,
                 page=0, total=0, notes=""):
    slide = blank_slide(prs)
    add_title_bar(slide, title, section=section, accent=accent)
    add_rect(slide, T.MARGIN_X, Inches(1.65), Inches(2.2), Inches(0.42),
             fill=accent)
    add_text(slide, T.MARGIN_X, Inches(1.66), Inches(2.2), Inches(0.4),
             "ANSWER", size=12, bold=True, color=T.BG_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, T.MARGIN_X, Inches(2.25), Inches(12.3), Inches(1.2),
             answer, size=24, bold=True, color=T.ACCENT_YELLOW)
    add_rect(slide, T.MARGIN_X, Inches(3.6), Inches(12.3), Inches(3.2),
             fill=T.BG_PANEL)
    add_text(slide, T.MARGIN_X + Inches(0.3), Inches(3.75),
             Inches(11.7), Inches(2.9), explanation,
             size=16, color=T.FG_TEXT)
    add_footer(slide, page=page, total=total)
    add_notes(slide, notes)
    return slide


def lab_callout_slide(prs, *, title: str, section: str, lab_name: str,
                      summary: str, bullets: Sequence[str],
                      accent=T.ACCENT_PURPLE, page=0, total=0, notes=""):
    slide = blank_slide(prs)
    add_title_bar(slide, title, section=section, accent=accent)
    # lab banner
    add_rect(slide, T.MARGIN_X, Inches(1.7), Inches(12.3), Inches(0.6),
             fill=accent)
    add_text(slide, T.MARGIN_X + Inches(0.3), Inches(1.7),
             Inches(12), Inches(0.6),
             f"LAB · {lab_name}", size=16, bold=True,
             color=T.BG_DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, T.MARGIN_X, Inches(2.45), Inches(12.3), Inches(0.8),
             summary, size=18, color=T.FG_MUTED, italic=True)
    add_bullets(slide, T.MARGIN_X, Inches(3.4), Inches(12.3),
                Inches(3.4), bullets, size=18,
                bullet_color=accent)
    add_footer(slide, page=page, total=total)
    add_notes(slide, notes)
    return slide
